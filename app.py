import io
import os
import subprocess
import tempfile
from datetime import datetime

import dropbox
import img2pdf
from flask import Flask, jsonify, request, send_file
from pypdf import PdfReader, PdfWriter
from pptx import Presentation

app = Flask(__name__)

# ── Environment variables (set these in Railway) ──────────────────────────────
API_SECRET            = os.environ.get("API_SECRET", "change-me-in-env")
DROPBOX_APP_KEY       = os.environ.get("DROPBOX_APP_KEY", "")
DROPBOX_APP_SECRET    = os.environ.get("DROPBOX_APP_SECRET", "")
DROPBOX_REFRESH_TOKEN = os.environ.get("DROPBOX_REFRESH_TOKEN", "")
DROPBOX_FILE_PATH     = os.environ.get("DROPBOX_FILE_PATH", "/EMCC_Certificate_TEMPLATE.pptx")


def download_template() -> bytes:
    """Download the PPTX template from private Dropbox using refresh token."""
    dbx = dropbox.Dropbox(
        oauth2_refresh_token=DROPBOX_REFRESH_TOKEN,
        app_key=DROPBOX_APP_KEY,
        app_secret=DROPBOX_APP_SECRET,
    )
    _, response = dbx.files_download(DROPBOX_FILE_PATH)
    return response.content


def replace_placeholders(prs: Presentation, replacements: dict) -> Presentation:
    """Replace {{PLACEHOLDER}} tokens in all text runs of a presentation."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for placeholder, value in replacements.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)
    return prs


def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    """
    Convert PPTX → PNG → PDF pipeline for pixel-perfect output.
    LibreOffice renders to PNG (preserving all fonts/layout),
    then img2pdf combines the PNG into a PDF without any re-encoding distortion.
    """
    profile_dir = os.path.join(output_dir, "soffice_profile")
    os.makedirs(profile_dir, exist_ok=True)

    # Step 1: PPTX → PNG via LibreOffice
    result = subprocess.run(
        [
            "soffice",
            "--headless",
            "--norestore",
            f"-env:UserInstallation=file://{profile_dir}",
            "--convert-to", "png",
            "--outdir", output_dir,
            pptx_path,
        ],
        capture_output=True,
        text=True,
        timeout=120,
        env={**os.environ, "HOME": "/tmp"},
    )

    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    png_path = os.path.join(output_dir, base_name + ".png")

    if not os.path.exists(png_path):
        stderr_clean = "\n".join(
            line for line in result.stderr.splitlines()
            if "javaldx" not in line and line.strip()
        )
        raise RuntimeError(
            f"LibreOffice PNG export failed.\n"
            f"Return code: {result.returncode}\n"
            f"Stderr: {stderr_clean}\n"
            f"Stdout: {result.stdout.strip()}"
        )

    #  Step 2: PNG → PDF via img2pdf (lossless, exact dimensions)
    pdf_path = os.path.join(output_dir, base_name + ".pdf")
    with open(pdf_path, "wb") as f:
        f.write(img2pdf.convert(png_path))

    return pdf_path


def set_pdf_metadata(pdf_path: str, first_name: str, last_name: str, issued_date: str) -> None:
    """Embed custom metadata into the generated PDF."""
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.add_metadata(reader.metadata or {})
    writer.add_metadata({
        "/Title":    f"EMCC USA Membership Certificate — {first_name} {last_name}",
        "/Author":   "EMCC USA",
        "/Subject":  "Individual Membership Certificate",
        "/Keywords": f"EMCC, membership, certificate, {first_name} {last_name}",
        "/Creator":  "EMCC USA Certificate Generator",
        "/Producer": "EMCC USA",
        "/IssuedTo": f"{first_name} {last_name}",
        "/IssuedOn": issued_date,
    })

    with open(pdf_path, "wb") as f:
        writer.write(f)


def format_date(dt: datetime) -> str:
    """Format a datetime as '27 February 2026'."""
    return dt.strftime("%B %-d, %Y")


@app.route("/generate-certificate", methods=["POST"])
def generate_certificate():
    # Authenticate
    if request.headers.get("X-API-Secret", "") != API_SECRET:
        return jsonify({"error": "Unauthorized"}), 401

    # Parse input
    body = request.get_json(force=True)
    data = body.get("data", body)  # unwrap "data" key if present
    first_name = data.get("first_name", "").strip()
    last_name  = data.get("last_name", "").strip()
    issued_date_str = data.get("issued_date") or data.get("paid_date")

    if not first_name or not last_name:
        return jsonify({"error": "first_name and last_name are required"}), 400

    # Compute dates
    # Strip timezone offset if present (e.g. "2026-03-01T00:00:00.000-05:00")
    issued_dt = None
    if issued_date_str:
        issued_dt = datetime.fromisoformat(issued_date_str)
    else:
        issued_dt = datetime.today()
    valid_dt = issued_dt.replace(year=issued_dt.year + 1)

    replacements = {
        "{{FIRST_NAME}}":  first_name,
        "{{LAST_NAME}}":   last_name,
        "{{ISSUED_DATE}}": format_date(issued_dt),
        "{{VALID_DATE}}":  format_date(valid_dt),
    }

    # Download template from private Dropbox
    try:
        template_bytes = download_template()
    except Exception as e:
        return jsonify({"error": f"Failed to download template: {str(e)}"}), 500

    # Generate certificate
    with tempfile.TemporaryDirectory() as tmpdir:
        prs = Presentation(io.BytesIO(template_bytes))
        prs = replace_placeholders(prs, replacements)

        pptx_out = os.path.join(tmpdir, f"certificate_{first_name}_{last_name}.pptx")
        prs.save(pptx_out)

        try:
            pdf_path = convert_pptx_to_pdf(pptx_out, tmpdir)
        except Exception as e:
            return jsonify({"error": f"PDF conversion failed: {str(e)}"}), 500

        try:
            set_pdf_metadata(pdf_path, first_name, last_name, format_date(issued_dt))
        except Exception:
            pass  # metadata is non-critical

        return send_file(
            pdf_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name=f"EMCC_USA_Certificate_{first_name}_{last_name}.pdf",
        )


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
